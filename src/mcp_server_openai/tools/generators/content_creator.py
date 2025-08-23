from __future__ import annotations

import re
from collections.abc import Sequence
from pathlib import Path
from typing import Any


def _safe_name(s: str) -> str:
    """
    Lowercase, replace non-alnum with '-', collapse repeats, and trim.
    """
    s = s.strip().lower()
    s = re.sub(r"[^a-z0-9._-]+", "-", s)
    s = re.sub(r"-+", "-", s).strip("-")
    return s or "untitled"


def _ensure_output_dir(client: str, project: str) -> str:
    """
    Ensure the output/<client>/<project> directory exists (portable).
    """
    base = Path("output")
    target = base / _safe_name(client) / _safe_name(project)
    target.mkdir(parents=True, exist_ok=True)
    return str(target)


def _heuristic_outline(brief: str, source_points: Sequence[str], num_slides: int) -> list[tuple[str, list[str]]]:
    """
    Tiny heuristic: slide 1 = Overview, slide N = Next steps, rest from source points.
    """
    slides: list[tuple[str, list[str]]] = []
    if num_slides < 1:
        num_slides = 1

    # Slide 1
    slides.append(("Overview", [brief]))

    # Middle slides from source points
    middle = max(0, num_slides - 2)
    for i in range(middle):
        chunk = source_points[i * 3 : (i + 1) * 3]
        title = f"Key updates {i + 1}"
        bullets = list(chunk) if chunk else ["TBD"]
        slides.append((title, bullets))

    # Final slide
    slides.append(("Next steps", ["Align on priorities", "Confirm owners & dates", "Plan follow-up"]))

    # Adjust if we over/undershot
    if len(slides) > num_slides:
        slides = slides[:num_slides]
    while len(slides) < num_slides:
        slides.append((f"Extra {len(slides)}", ["TBD"]))

    return slides


def _create_ppt_from_outline(outline: Sequence[tuple[str, Sequence[str]]]) -> Any:
    """
    Build a Presentation from an outline. We avoid type annotations to keep mypy
    happy without python-pptx stubs; return type is 'Any'.
    """
    from pptx import Presentation  # runtime import

    prs = Presentation()
    for title, bullets in outline:
        layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = str(title)

        body = slide.placeholders[1].text_frame
        # clear any default paragraph
        if body.paragraphs:
            body.clear()

        first = True
        for b in bullets:
            text = str(b)
            if first:
                body.text = text
                first = False
            else:
                p = body.add_paragraph()
                p.text = text
                p.level = 0

    return prs


def _save_ppt(prs: Any, client: str, project: str) -> str:
    """
    Save the presentation under output/<client>/<project>/content.pptx.
    Ensures directory exists even if caller monkey-patches _ensure_output_dir.
    """
    out_dir = Path(_ensure_output_dir(client, project))
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "content.pptx"
    prs.save(str(out_path))
    return str(out_path)


# Public tool entry (kept simple; tests call internals)
def create_content_ppt(
    client_name: str,
    project_name: str,
    source_content_type: str,
    source_content_details: Sequence[str],
    target_content_type: str,
    number_of_slides: int,
    content_brief: str,
) -> dict[str, Any]:
    """
    Minimal entrypoint used by MCP; returns metadata about the generated deck.
    """
    outline = _heuristic_outline(content_brief, list(source_content_details), number_of_slides)
    prs = _create_ppt_from_outline(outline)
    path = _save_ppt(prs, client_name, project_name)
    return {
        "path": path,
        "slides": getattr(prs, "slides", []).__len__(),
        "client": client_name,
        "project": project_name,
        "target": target_content_type,
    }
