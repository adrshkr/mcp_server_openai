"""
Content Creator tool (Milestone 2, PPT-first).

Generates a PPTX from:
- Client Name (S3 folder name later)
- Project Name (subfolder)
- Source Content:
    - Type: "Highlight" or "Temp_Resource"
    - Details: list[str] (for now we accept strings only; python data to be supported later)
- Target Content Type: "PPT" (only PPT supported in this milestone)
- Target Content Specs:
    - Number of Slides: int
- Content Brief: 1-2 paragraph abstract
- LLM config placeholders (primary/secondary model and keys) accepted but not used yet.
"""
from __future__ import annotations

import os
import re
from dataclasses import dataclass
from typing import TYPE_CHECKING, List, Optional, Literal

from pptx import Presentation
from pptx.util import Pt

if TYPE_CHECKING:
  from mcp.server.fastmcp import FastMCP


@dataclass
class LLMConfig:
  primary_model: Optional[str] = None
  secondary_model: Optional[str] = None
  primary_api_key: Optional[str] = None
  secondary_api_key: Optional[str] = None


def _safe_name(s: str) -> str:
  s = s.strip().replace(" ", "_")
  s = re.sub(r"[^A-Za-z0-9_.-]", "", s)
  return s or "untitled"


def _ensure_output_dir(client: str, project: str) -> str:
  base = os.path.join("output", _safe_name(client), _safe_name(project))
  os.makedirs(base, exist_ok=True)
  return base


def _build_slide(prs: Presentation, title: str, bullets: List[str]) -> None:
  layout = prs.slide_layouts[1]  # Title and Content
  slide = prs.slides.add_slide(layout)
  slide.shapes.title.text = title
  body = slide.shapes.placeholders[1].text_frame
  body.clear()
  for i, b in enumerate(bullets):
    if i == 0:
      body.text = b
    else:
      body.add_paragraph().text = b
  for p in body.paragraphs:
    for run in p.runs:
      run.font.size = Pt(18)


def _heuristic_outline(brief: str, source_lines: List[str], num_slides: int) -> List[tuple[str, List[str]]]:
  slides: List[tuple[str, List[str]]] = []
  intro = ("Overview", [brief[:200] + ("..." if len(brief) > 200 else "")])
  slides.append(intro)
  remaining = max(0, num_slides - 1)
  chunks = [source_lines[i:i+5] for i in range(0, len(source_lines), 5)]
  for idx in range(remaining):
    block = chunks[idx] if idx < len(chunks) else []
    if not block:
      block = ["(placeholder)"]
    slides.append((f"Section {idx+1}", block))
  return slides[:num_slides]


def _create_ppt_from_outline(outline: List[tuple[str, List[str]]]) -> Presentation:
  prs = Presentation()
  title_layout = prs.slide_layouts[0]
  title_slide = prs.slides.add_slide(title_layout)
  title_slide.shapes.title.text = "Generated Deck"
  if len(title_slide.shapes.placeholders) > 1:
    title_slide.placeholders[1].text = "via mcp_server_openai"
  for title, bullets in outline:
    _build_slide(prs, title, bullets)
  return prs


def _save_ppt(prs: Presentation, client: str, project: str) -> str:
  out_dir = _ensure_output_dir(client, project)
  # Be defensive: ensure dir exists even if _ensure_output_dir was monkeypatched
  os.makedirs(out_dir, exist_ok=True)
  filename = "content.pptx"
  path = os.path.join(out_dir, filename)
  prs.save(path)
  return path


def register(mcp: "FastMCP") -> None:
  @mcp.tool(
    name="content.create",
    description="Create content from a brief and source materials. PPT-only in this milestone."
  )
  async def create_content(
    client_name: str,
    project_name: str,
    source_content_type: Literal["Highlight", "Temp_Resource"],
    source_content_details: List[str],
    target_content_type: Literal["PPT"],
    number_of_slides: int,
    content_brief: str,
    primary_llm: Optional[str] = None,
    secondary_llm: Optional[str] = None,
    primary_api_key: Optional[str] = None,
    secondary_api_key: Optional[str] = None,
  ) -> dict[str, str | int]:
    _ = LLMConfig(primary_llm, secondary_llm, primary_api_key, secondary_api_key)
    source_lines = [s.strip() for s in (source_content_details or []) if s and s.strip()]
    ns = max(1, int(number_of_slides))
    outline = _heuristic_outline(content_brief or "", source_lines, ns)
    prs = _create_ppt_from_outline(outline)
    path = _save_ppt(prs, client_name, project_name)
    return {
      "result": "ok",
      "output_type": "PPT",
      "slides": len(prs.slides),
      "path": path,
    }